"""CAPACITY CONNECT — SIH 2026 Idea Submission, official SIH template format (6 slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__),
                   "SIH2026_IDEA_CAPACITY_CONNECT_26075.pptx")

# ---- SIH template palette (light theme) ----
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x1E, 0x6F, 0xBA)
LTBLUE = RGBColor(0xE8, 0xF1, 0xFA)
ORANGE = RGBColor(0xF2, 0x7C, 0x1E)
LTORANGE = RGBColor(0xFD, 0xF0, 0xE2)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
LTGREEN = RGBColor(0xE9, 0xF5, 0xEE)
GREY = RGBColor(0x44, 0x4C, 0x57)
LTGREY = RGBColor(0xF3, 0xF5, 0xF8)
BORDER = RGBColor(0xC9, 0xD5, 0xE3)
RED = RGBColor(0xC0, 0x39, 0x2B)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

TEAM = "Team Name"          # <- edit
TEAM_ID = "Team ID"         # <- edit


def rect(s, x, y, w, h, fill=WHITE, line=None, radius=False, lw=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=3, line_spacing=1.08):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # Convention: a list  -> its own paragraph containing several inline runs;
    #             a tuple -> a single-run paragraph.
    # A flat list of tuples is treated as ONE paragraph of inline runs.
    if runs and all(isinstance(i, (tuple, str)) for i in runs):
        groups = [list(runs)]
    else:
        groups = [g if isinstance(g, list) else [g] for g in runs]
    first = True
    for grp in groups:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for item in grp:
            if isinstance(item, str):
                item = (item, 11, GREY, False)
            txt, size, color, bold = (list(item) + [False])[:4]
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def slide(title=None, num=None):
    s = prs.slides.add_slide(BLANK)
    bg = rect(s, 0, 0, W, H, WHITE)
    # left accent spine
    rect(s, 0, 0, Inches(0.14), H, NAVY)
    rect(s, 0, 0, Inches(0.14), Inches(2.2), ORANGE)
    if title:
        text(s, Inches(0.5), Inches(0.3), Inches(10.5), Inches(0.6),
             [(title, 30, NAVY, True)])
        rect(s, Inches(0.5), Inches(0.92), Inches(2.0), Pt(4), ORANGE)
        rect(s, Inches(2.5), Inches(0.92), Inches(10.0), Pt(4), BORDER)
    # footer
    rect(s, 0, H - Inches(0.42), W, Inches(0.42), LTGREY)
    text(s, Inches(0.5), H - Inches(0.31), Inches(4.5), Inches(0.25),
         [("@SIH Idea submission- Template", 10, GREY, False)])
    if num:
        text(s, Inches(6.3), H - Inches(0.31), Inches(0.7), Inches(0.25),
             [(str(num), 10, NAVY, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(8.3), H - Inches(0.31), Inches(4.5), Inches(0.25),
         [(TEAM, 10, GREY, True)], align=PP_ALIGN.RIGHT)
    return s


# --- text measurement (DejaVu metrics scaled to approximate Calibri) ---
from PIL import Image, ImageDraw, ImageFont
_M = ImageDraw.Draw(Image.new("RGB", (10, 10)))
_FCACHE = {}
CALIBRI_FACTOR = 0.97   # Calibri is ~10% narrower than DejaVu Sans


def _f(size_pt, bold=False):
    k = (round(size_pt, 1), bold)
    if k not in _FCACHE:
        path = "/usr/share/fonts/truetype/dejavu/DejaVuSans%s.ttf" % ("-Bold" if bold else "")
        _FCACHE[k] = ImageFont.truetype(path, max(6, int(round(size_pt * 4))))
    return _FCACHE[k]


def est_lines(txt, size_pt, width_in, bold=False):
    """Conservative line count for `txt` at `size_pt` inside `width_in` inches."""
    f = _f(size_pt, bold)
    w_px = _M.textlength(txt, font=f) / 4.0          # back to points
    w_pt = w_px * CALIBRI_FACTOR
    avail_pt = width_in * 72.0
    return max(1, int(-(-w_pt // avail_pt)))


def bullet_block(s, x, y, w, items, size=11.5, gap=0.32, color=GREY,
                 dot=ORANGE, bold_head=True):
    """items: str, or (head, tail)."""
    cy = y
    for it in items:
        if it == "":
            cy += Inches(gap * 0.4)
            continue
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, cy + Inches(0.055), Pt(5.5), Pt(5.5))
        d.fill.solid(); d.fill.fore_color.rgb = dot
        d.line.fill.background(); d.shadow.inherit = False
        if isinstance(it, tuple):
            runs = [(it[0] + " ", size, NAVY, bold_head), (it[1], size, color, False)]
        else:
            runs = [(it, size, color, False)]
        text(s, x + Inches(0.2), cy, w - Inches(0.2), Inches(0.4), runs,
             line_spacing=1.12)
        raw = (it[0] + " " + it[1]) if isinstance(it, tuple) else it
        lines = est_lines(raw, size, (w - Inches(0.2)) / Inches(1))
        line_h = size * 1.12 / 72.0
        cy += Inches(max(gap, line_h * lines + 0.19))
    return cy


def section(s, x, y, w, label, color=BLUE, fill=LTBLUE):
    b = rect(s, x, y, w, Inches(0.38), fill, line=None, radius=True)
    text(s, x + Inches(0.16), y + Inches(0.07), w - Inches(0.3), Inches(0.28),
         [(label, 12, color, True)])
    return y + Inches(0.52)


def flow(s, x, y, w, labels, color=BLUE, fill=LTBLUE, size=9, h=0.55):
    n = len(labels)
    gap = Inches(0.16)
    bw = (w - gap * (n - 1)) / n
    for i, lab in enumerate(labels):
        bx = x + (bw + gap) * i
        rect(s, bx, y, bw, Inches(h), fill, line=color, radius=True)
        text(s, bx + Inches(0.04), y, bw - Inches(0.08), Inches(h),
             [(lab, size, NAVY, True)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        if i < n - 1:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, bx + bw + Inches(0.015),
                                   y + Inches(h / 2 - 0.055), Inches(0.13), Inches(0.11))
            a.fill.solid(); a.fill.fore_color.rgb = ORANGE
            a.line.fill.background(); a.shadow.inherit = False


# ==================== SLIDE 1 — TITLE PAGE ====================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, W, Inches(0.16), ORANGE)
rect(s, 0, Inches(0.16), W, Inches(0.06), NAVY)
rect(s, 0, H - Inches(0.42), W, Inches(0.42), LTGREY)
text(s, Inches(0.5), H - Inches(0.31), Inches(5), Inches(0.25),
     [("@SIH Idea submission- Template", 10, GREY, False)])
text(s, Inches(6.3), H - Inches(0.31), Inches(0.7), Inches(0.25),
     [("1", 10, NAVY, True)], align=PP_ALIGN.CENTER)
text(s, Inches(8.3), H - Inches(0.31), Inches(4.5), Inches(0.25),
     [(TEAM, 10, GREY, True)], align=PP_ALIGN.RIGHT)

text(s, Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.5),
     [("TITLE PAGE", 20, ORANGE, True)])
text(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.7),
     [("SMART INDIA HACKATHON 2026", 36, NAVY, True)])
rect(s, Inches(0.8), Inches(2.0), Inches(2.6), Pt(5), ORANGE)

fields = [("Problem Statement ID", "26075"),
          ("Problem Statement Title",
           "CAPACITY CONNECT — A Digital Capacity Building and Learning Management Portal"),
          ("Theme", "Smart Education"),
          ("PS Category", "Software"),
          ("Team ID", TEAM_ID),
          ("Team Name (Registered on portal)", TEAM)]
y = Inches(2.35)
for i, (k, v) in enumerate(fields):
    h = Inches(0.62)
    rect(s, Inches(0.8), y, Inches(8.3), h, LTGREY if i % 2 == 0 else WHITE,
         line=BORDER)
    rect(s, Inches(0.8), y, Pt(4), h, ORANGE if i < 2 else BLUE)
    text(s, Inches(1.0), y + Inches(0.09), Inches(3.1), Inches(0.45),
         [(k, 12, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.2), y + Inches(0.05), Inches(4.75), Inches(0.52),
         [(v, 12, GREY, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y += h + Inches(0.09)

rect(s, Inches(9.45), Inches(2.35), Inches(3.1), Inches(4.05), LTBLUE, line=BORDER)
text(s, Inches(9.7), Inches(2.6), Inches(2.6), Inches(0.3),
     [("ORGANISATION", 11, BLUE, True)])
text(s, Inches(9.7), Inches(2.95), Inches(2.6), Inches(0.8),
     [("Ministry of Earth Sciences (MoES)", 13, NAVY, True)], line_spacing=1.15)
text(s, Inches(9.7), Inches(3.72), Inches(2.6), Inches(0.3),
     [("DEPARTMENT", 11, BLUE, True)])
text(s, Inches(9.7), Inches(4.05), Inches(2.6), Inches(0.7),
     [("India Meteorological Department", 13, NAVY, True)], line_spacing=1.15)
rect(s, Inches(9.7), Inches(4.85), Inches(2.6), Pt(3), BORDER)
text(s, Inches(9.7), Inches(5.05), Inches(2.6), Inches(1.2),
     [("Offline-first, secure, role-based capacity building for field, vessel and "
       "observatory personnel.", 11, GREY, False)], line_spacing=1.25)

# ==================== SLIDE 2 — IDEA TITLE / PROPOSED SOLUTION ====================
s = slide("IDEA TITLE", 2)
text(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.35),
     [("CAPACITY CONNECT — one offline-first portal for learning, assessment, "
       "competency mapping and AI-assisted field troubleshooting", 14, ORANGE, True)])

y = section(s, Inches(0.5), Inches(1.58), Inches(7.6),
            "❖ Proposed Solution (Describe your Idea/Solution/Prototype)")
text(s, Inches(0.5), y, Inches(7.6), Inches(0.3),
     [("Detailed explanation of the proposed solution", 11.5, NAVY, True)])
y += Inches(0.3)
y = bullet_block(s, Inches(0.62), y, Inches(7.4), [
    ("Single portal, three roles.", "Trainee, Trainer and Admin — secure signup/login with Admin approval and role management."),
    ("Trainee:", "professional profile (qualifications, experience, skills, interests, certificates), course enrollment, learning resources, subject-wise MCQ assessments and course feedback."),
    ("Trainer:", "competency declaration, questionnaires with deadlines, participation and performance monitoring, and a Trainer Library of recorded lectures, presentations and study material."),
    ("Admin:", "user approval, role management, dashboards for courses / enrollments / certifications / assessments / participation, plus homepage publishing of notifications, announcements and achievements."),
    ("Competency mapping engine", "scores trainers against subject requirements to objectively identify the right trainer."),
    ("Offline-first edge app:", "a real local database plus a local AI node on 127.0.0.1, so work continues at sea and at remote sites."),
], size=10, gap=0.28)

y2 = section(s, Inches(0.5), y + Inches(0.05), Inches(7.6), "How it addresses the problem",
             ORANGE, LTORANGE)
bullet_block(s, Inches(0.62), y2, Inches(7.4), [
    "Replaces scattered drives, mail and spreadsheets with one auditable source of truth for learning and competency.",
    "Removes the connectivity barrier: every feature works offline and reconciles automatically on reconnect.",
    "Turns tacit instrument know-how into retrievable, guided diagnostic procedures.",
], size=10, gap=0.28, dot=ORANGE)

rect(s, Inches(8.35), Inches(1.58), Inches(4.5), Inches(5.2), LTBLUE, line=BORDER)
text(s, Inches(8.6), Inches(1.75), Inches(4.0), Inches(0.3),
     [("Innovation and uniqueness of the solution", 12, BLUE, True)])
inno = [("Offline is the default, not a fallback.",
         "The local SQLite/Drift store is a functioning operating database, not a cache."),
        ("Local AI engine with no cloud LLM.",
         "Deterministic keyword→phrase→fuzzy→retrieval pipeline returns structured JSON on 127.0.0.1."),
        ("Learning fused with operations.",
         "A described fault maps to a real asset and to the exact training module that covers it."),
        ("3D Digital Twin guidance.",
         "SONAR-001 → Mesh_042 highlighted in CRITICAL state from locally cached GLTF/GLB."),
        ("Objective competency mapping.",
         "Trainer-to-subject fit is computed from skills, certificates and assessment outcomes."),
        ("Auditable sync ledger.",
         "Every offline action is a versioned transaction with explicit conflict handling.")]
cy = Inches(2.2)
for h_, d in inno:
    text(s, Inches(8.6), cy, Inches(4.0), Inches(0.3), [(h_, 11, NAVY, True)])
    text(s, Inches(8.6), cy + Inches(0.22), Inches(4.0), Inches(0.5),
         [(d, 10, GREY, False)], line_spacing=1.12)
    cy += Inches(0.75)

# ==================== SLIDE 3 — TECHNICAL APPROACH ====================
s = slide("TECHNICAL APPROACH", 3)
y = section(s, Inches(0.5), Inches(1.15), Inches(12.35),
            "Technologies to be used (programming languages, frameworks, hardware)")
stacks = [("CLIENT", BLUE, LTBLUE,
           ["Flutter 3 / Dart, Material 3", "Riverpod state · GoRouter",
            "Drift + SQLite (SQLCipher)", "flutter_secure_storage · Dio",
            "Android · Windows · Linux · Web"]),
          ("LOCAL AI", ORANGE, LTORANGE,
           ["Python 3.11 · FastAPI on 127.0.0.1", "Normalize · language detect · tokenize",
            "Keyword / phrase / RapidFuzz match", "TF-IDF knowledge retrieval",
            "Pluggable local LLM slot"]),
          ("3D & MEDIA", GREEN, LTGREEN,
           ["GLTF / GLB cached locally", "Mesh mapping + fault state layer",
            "Chunked resumable downloads", "Streaming video · PDF viewer",
            "Optional speech-to-text input"]),
          ("BACKEND & OPS", NAVY, LTGREY,
           ["FastAPI · SQLAlchemy · Pydantic", "PostgreSQL + Alembic migrations",
            "Celery workers · JWT + refresh", "Docker · GitHub Actions CI",
            "pytest · flutter test/analyze"])]
for i, (t, c, f, items) in enumerate(stacks):
    x = Inches(0.5 + i * 3.12)
    rect(s, x, y, Inches(2.95), Inches(1.72), f, line=BORDER)
    rect(s, x, y, Inches(2.95), Pt(4), c)
    text(s, x + Inches(0.15), y + Inches(0.13), Inches(2.6), Inches(0.28),
         [(t, 11.5, c, True)])
    for j, it in enumerate(items):
        text(s, x + Inches(0.15), y + Inches(0.45) + Inches(0.25) * j, Inches(2.7),
             Inches(0.24), [("· ", 10, c, True), (it, 9.5, GREY, False)])

y = y + Inches(1.95)
y = section(s, Inches(0.5), y, Inches(12.35),
            "Methodology and process for implementation", ORANGE, LTORANGE)

text(s, Inches(0.5), y, Inches(12.3), Inches(0.25),
     [("A. System architecture — edge application works with zero connectivity; cloud platform is independently deployable",
       10.5, NAVY, True)])
y += Inches(0.3)
flow(s, Inches(0.5), y, Inches(12.35),
     ["Flutter Client\n(Material 3)", "Local AI Node\n127.0.0.1", "Local DB\nSQLCipher",
      "Digital Twin\nGLTF/GLB", "Sync Ledger\n(encrypted)", "FastAPI\n/api/v1",
      "PostgreSQL\n+ workers", "Admin Portal\n& Dashboards"], size=8.5, h=0.58)

y += Inches(0.78)
text(s, Inches(0.5), y, Inches(12.3), Inches(0.25),
     [("B. Local AI troubleshooting pipeline — returns structured JSON, confidence is a defined weighted score",
       10.5, NAVY, True)])
y += Inches(0.3)
flow(s, Inches(0.5), y, Inches(12.35),
     ["Input", "Normalize", "Tokenize", "Keyword +\nPhrase", "Fuzzy\nMatch",
      "Knowledge\nRetrieval", "Component\n+ Severity", "3D Mesh\nMapping"],
     color=ORANGE, fill=LTORANGE, size=8.5, h=0.55)

y += Inches(0.72)
rect(s, Inches(0.5), y, Inches(12.35), Inches(0.72), LTGREY, line=BORDER)
text(s, Inches(0.7), y + Inches(0.09), Inches(12.0), Inches(0.55),
     [('Worked example: "Sonar transducer showing abnormal vibration and casing fracture" → ',
       10.5, GREY, False),
      ("SONAR-001 · Sonar Transducer Array · Mesh_042 · severity HIGH · confidence 0.94 → "
       "Digital Twin highlight + diagnostic procedure + diagnostic record saved offline.",
       10.5, NAVY, True)], line_spacing=1.2)

# ==================== SLIDE 4 — FEASIBILITY AND VIABILITY ====================
s = slide("FEASIBILITY AND VIABILITY", 4)
y = section(s, Inches(0.5), Inches(1.15), Inches(12.35), "Analysis of the feasibility of the idea")
feas = [("Technically proven stack",
         "Flutter, FastAPI and PostgreSQL are mature and fully open-source; one Dart codebase covers Android, Windows, Linux and Web."),
        ("No exotic hardware",
         "Runs on existing MoES/IMD field tablets and vessel workstations; backend is containerised on existing infrastructure."),
        ("Zero licensing cost",
         "Entirely open-source; no per-seat fees, so scaling to more users and institutes costs only storage and compute."),
        ("Incremental rollout",
         "Portal features usable from day one; AI and Digital Twin layers extend the same schema without rework.")]
cy = y
for h_, d in feas:
    col = 0 if feas.index((h_, d)) % 2 == 0 else 1
    pass
for i, (h_, d) in enumerate(feas):
    x = Inches(0.5 + (i % 2) * 6.25)
    yy = y + Inches((i // 2) * 0.76)
    rect(s, x, yy, Inches(6.0), Inches(0.68), LTGREEN, line=BORDER)
    rect(s, x, yy, Pt(4), Inches(0.68), GREEN)
    text(s, x + Inches(0.18), yy + Inches(0.08), Inches(5.6), Inches(0.25),
         [(h_, 11, NAVY, True)])
    text(s, x + Inches(0.18), yy + Inches(0.32), Inches(5.65), Inches(0.35),
         [(d, 9.5, GREY, False)], line_spacing=1.1)

y = y + Inches(1.62)
y = section(s, Inches(0.5), y, Inches(12.35),
            "Potential challenges and risks  →  Strategies for overcoming these challenges",
            ORANGE, LTORANGE)

rows = [("Poor / no connectivity at sea and remote sites",
         "Offline-first architecture: local operating DB + local AI node; sync ledger drains automatically on reconnect."),
        ("Large 3D models on low-end field devices",
         "Decimated LOD meshes by default, full-resolution GLB as optional download; renderer degrades to metadata-only view."),
        ("Offline credential caching weakens authentication",
         "Admin-approved devices only, Argon2id verifier (never the password), bounded grace window, audit + re-verification on sync."),
        ("Sync conflicts on shared records",
         "Per-entity strategies (version compare, field merge, manual queue); diagnostic records are never silently overwritten."),
        ("Digital literacy and user adoption",
         "Role-tailored simple UI, localisation-ready strings for Indian languages, optional voice input, in-app guided demo scenario."),
        ("Storage exhaustion on field devices",
         "Storage manager with per-category usage, resumable downloads, permissioned removal of optional assets only.")]
rect(s, Inches(0.5), y, Inches(12.35), Inches(0.34), NAVY)
text(s, Inches(0.7), y + Inches(0.06), Inches(5.0), Inches(0.25),
     [("CHALLENGE / RISK", 10.5, WHITE, True)])
text(s, Inches(6.2), y + Inches(0.06), Inches(6.4), Inches(0.25),
     [("MITIGATION STRATEGY", 10.5, WHITE, True)])
for i, (a, b) in enumerate(rows):
    ry = y + Inches(0.34 + i * 0.47)
    rect(s, Inches(0.5), ry, Inches(12.35), Inches(0.47),
         LTGREY if i % 2 == 0 else WHITE, line=BORDER)
    text(s, Inches(0.7), ry + Inches(0.07), Inches(5.2), Inches(0.4),
         [(a, 10, NAVY, True)], line_spacing=1.08)
    text(s, Inches(6.2), ry + Inches(0.07), Inches(6.45), Inches(0.4),
         [(b, 9.5, GREY, False)], line_spacing=1.08)

# ==================== SLIDE 5 — IMPACT AND BENEFITS ====================
s = slide("IMPACT AND BENEFITS", 5)
y = section(s, Inches(0.5), Inches(1.15), Inches(12.35), "Potential impact on the target audience")
imp = [("TRAINEES / FIELD ENGINEERS", BLUE, LTBLUE,
        ["Learn, test and certify with zero connectivity — at sea and at remote sites",
         "Guided AI diagnostics instead of waiting for shore support",
         "Portable, verifiable competency profile"]),
       ("TRAINERS", ORANGE, LTORANGE,
        ["One versioned library for lectures, presentations and material",
         "Questionnaires with deadlines, auto-scoring, live participation view",
         "Feedback shows which content actually works"]),
       ("ADMINISTRATION / MoES", GREEN, LTGREEN,
        ["Live dashboards of enrollment, certification and assessment",
         "Objective trainer selection via competency mapping",
         "Complete audit trail for governance and reporting"])]
for i, (t, c, f, items) in enumerate(imp):
    x = Inches(0.5 + i * 4.15)
    rect(s, x, y, Inches(3.95), Inches(1.92), f, line=BORDER)
    rect(s, x, y, Inches(3.95), Pt(4), c)
    text(s, x + Inches(0.16), y + Inches(0.14), Inches(3.6), Inches(0.28),
         [(t, 11, c, True)])
    for j, it in enumerate(items):
        text(s, x + Inches(0.16), y + Inches(0.46) + Inches(0.47) * j, Inches(3.6),
             Inches(0.44), [("· ", 9.5, c, True), (it, 9.5, GREY, False)],
             line_spacing=1.1)

y = y + Inches(2.14)
y = section(s, Inches(0.5), y, Inches(12.35),
            "Benefits of the solution (social, economic, environmental)", ORANGE, LTORANGE)
ben = [("SOCIAL", "Equitable access to training for remote and shipboard staff; skill "
                  "recognition through verifiable certificates; localisation-ready for "
                  "Indian-language expansion; safer field work through guided procedures."),
       ("ECONOMIC", "Lower instrument downtime and fewer avoidable failures; reduced travel "
                    "and physical-classroom cost; open-source stack with no licensing fees; "
                    "reusable across INCOIS, NCPOR and NIOT on the same core."),
       ("ENVIRONMENTAL", "Fewer service trips and paper-free assessments; longer instrument "
                         "life through timely preventive maintenance; better-maintained ocean "
                         "and weather sensors mean higher-quality earth-science data.")]
for i, (t, d) in enumerate(ben):
    x = Inches(0.5 + i * 4.15)
    rect(s, x, y, Inches(3.95), Inches(1.35), WHITE, line=BORDER)
    rect(s, x, y, Pt(4), Inches(1.35), NAVY)
    text(s, x + Inches(0.18), y + Inches(0.12), Inches(3.6), Inches(0.25),
         [(t, 11, NAVY, True)])
    text(s, x + Inches(0.18), y + Inches(0.42), Inches(3.6), Inches(0.85),
         [(d, 9.5, GREY, False)], line_spacing=1.15)

y = y + Inches(1.55)
rect(s, Inches(0.5), y, Inches(12.35), Inches(0.62), LTGREY, line=BORDER)
metrics = [("100%", "core features usable offline"), ("↓ 40%", "target fault-resolution time"),
           ("3", "roles, one governed platform"), ("4", "platforms, one codebase")]
for i, (k, v) in enumerate(metrics):
    x = Inches(0.6 + i * 3.07)
    text(s, x, y + Inches(0.1), Inches(1.3), Inches(0.3), [(k, 15, ORANGE, True)])
    text(s, x + Inches(1.35), y + Inches(0.18), Inches(1.7), Inches(0.3),
         [(v, 10, GREY, False)])

# ==================== SLIDE 6 — RESEARCH AND REFERENCES ====================
s = slide("RESEARCH AND REFERENCES", 6)
y = section(s, Inches(0.5), Inches(1.15), Inches(12.35),
            "Details / Links of the reference and research work")

groups = [
    ("Problem & organisational context", BLUE, [
        "SIH 2026 Problem Statement 26075 — Ministry of Earth Sciences (MoES) — sih.gov.in",
        "moes.gov.in  |  India Meteorological Department — mausam.imd.gov.in",
        "MoES Annual Report — capacity building & human-resource development",
        "National Education Policy 2020 — digital, competency-based learning",
        "SIH 2026 portal guidelines — idea submission format and evaluation criteria",
    ]),
    ("Standards & frameworks referenced", ORANGE, [
        "SCORM / xAPI (Experience API) — learning-record interoperability",
        "IEEE 1484 LOM — learning-object metadata for the course catalog",
        "NIST SP 800-63B — digital identity & authentication guidance",
        "OWASP ASVS + Mobile Top 10 — application security verification baseline",
        "NIST SP 800-38D — AES-256-GCM authenticated encryption",
        "ISO/IEC 27001 Annex A — access control, logging, audit-trail controls",
    ]),
    ("Technology documentation", GREEN, [
        "Flutter & Dart — docs.flutter.dev | Riverpod, GoRouter, Drift",
        "FastAPI — fastapi.tiangolo.com  |  SQLAlchemy, Alembic, Pydantic official docs",
        "SQLCipher — full-database encryption for SQLite (zetetic.net/sqlcipher)",
        "Khronos glTF 2.0 specification — 3D asset format for the Digital Twin",
        "PostgreSQL 16 docs | Docker Compose deployment reference",
    ]),
    ("Domain & prior art studied", NAVY, [
        "Digital Twin concepts for marine instrumentation (ISO 23247 framing)",
        "Argo programme documentation — argo.ucsd.edu — float operations",
        "Offline-first patterns; CRDT vs. ledger-based synchronisation literature",
        "Open-source LMS study (Moodle, Open edX) — role model & gap analysis",
        "RapidFuzz / TF-IDF retrieval literature — basis of the confidence score",
    ]),
]
cy = y
for i, (t, c, items) in enumerate(groups):
    x = Inches(0.5 + (i % 2) * 6.25)
    yy = y + Inches((i // 2) * 2.45)
    h = Inches(2.3)
    rect(s, x, yy, Inches(6.0), h, WHITE, line=BORDER)
    rect(s, x, yy, Inches(6.0), Pt(4), c)
    text(s, x + Inches(0.18), yy + Inches(0.14), Inches(5.6), Inches(0.28),
         [(t, 11.5, c, True)])
    for j, it in enumerate(items):
        text(s, x + Inches(0.18), yy + Inches(0.48) + Inches(0.285) * j, Inches(5.6),
             Inches(0.27), [("· ", 9, c, True), (it, 9, GREY, False)],
             line_spacing=1.05)

text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.3),
     [("Note: ", 9.5, NAVY, True),
      ("all cited standards and libraries are open or publicly available; the solution uses no "
       "proprietary or licence-restricted component.", 9.5, GREY, False)])

prs.save(OUT)
print("saved", OUT)
