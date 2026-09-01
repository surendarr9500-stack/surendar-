"""Lightweight PPTX -> PNG previewer for the decks in this folder.

Only supports the shape vocabulary used by our generator scripts
(rectangles, rounded rectangles, ovals, right arrows, text boxes with runs).
Used for visual verification, not as a general-purpose renderer.
"""
import os
import sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

FONT_DIR = "/usr/share/fonts/truetype/dejavu"
REG = os.path.join(FONT_DIR, "DejaVuSans.ttf")
BOLD = os.path.join(FONT_DIR, "DejaVuSans-Bold.ttf")
MONO = os.path.join(FONT_DIR, "DejaVuSansMono.ttf")
_cache = {}


def font(size_pt, bold=False, mono=False, scale=1.0):
    px = max(6, int(round(size_pt * scale * 96 / 72)))
    key = (px, bold, mono)
    if key not in _cache:
        path = MONO if mono else (BOLD if bold else REG)
        _cache[key] = ImageFont.truetype(path, px)
    return _cache[key]


def rgb(c, default=(0, 0, 0)):
    try:
        if c is None:
            return default
        return (c[0], c[1], c[2])
    except Exception:
        return default


def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is None:
            return None
        if f.type == 1:  # solid
            return rgb(f.fore_color.rgb)
    except Exception:
        pass
    return None


def shape_line(sh):
    try:
        lf = sh.line.fill
        if lf.type == 1:
            w = sh.line.width
            return rgb(sh.line.color.rgb), max(1, int((w or 12700) / 12700 * 96 / 72))
    except Exception:
        pass
    return None, 0


def wrap(draw, txt, fnt, maxw):
    if "\n" in txt:
        out = []
        for seg in txt.split("\n"):
            out.extend(wrap(draw, seg, fnt, maxw))
        return out
    words = txt.split(" ")
    lines, cur = [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if draw.textlength(t, font=fnt) <= maxw or not cur:
            cur = t
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def render(path, outdir, scale=1.0):
    prs = Presentation(path)
    SW, SH = prs.slide_width, prs.slide_height
    px = lambda v: int(v / 914400 * 96 * scale)
    W, H = px(SW), px(SH)
    os.makedirs(outdir, exist_ok=True)
    paths = []
    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(img)
        for sh in slide.shapes:
            if sh.left is None:
                continue
            x0, y0 = px(sh.left), px(sh.top)
            x1, y1 = px(sh.left + sh.width), px(sh.top + sh.height)
            if x1 < x0:
                x0, x1 = x1, x0
            if y1 < y0:
                y0, y1 = y1, y0
            name = sh.shape_type
            if sh.has_text_frame and str(sh.shape_type).startswith("TEXT_BOX"):
                pass
            else:
                fill = shape_fill(sh)
                lcol, lw = shape_line(sh)
                st = str(sh.name)
                if fill or lcol:
                    if "Oval" in st:
                        d.ellipse([x0, y0, x1, y1], fill=fill, outline=lcol, width=lw or 1)
                    elif "Arrow" in st:
                        mid = (y0 + y1) / 2
                        d.polygon([(x0, y0), (x1 - (x1 - x0) * 0.4, y0),
                                   (x1, mid), (x1 - (x1 - x0) * 0.4, y1), (x0, y1)],
                                  fill=fill or lcol)
                    elif "Rounded" in st:
                        r = max(2, int(min(x1 - x0, y1 - y0) * 0.12))
                        d.rounded_rectangle([x0, y0, x1, y1], radius=r, fill=fill,
                                            outline=lcol, width=lw or 1)
                    else:
                        d.rectangle([x0, y0, x1, y1], fill=fill, outline=lcol, width=lw or 1)
            # text
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            if not tf.text.strip():
                continue
            paras = []
            for p in tf.paragraphs:
                runs = [(r.text, (r.font.size.pt if r.font.size else 11),
                         rgb(r.font.color.rgb if r.font.color and r.font.color.type is not None else None,
                             (40, 40, 40)),
                         bool(r.font.bold),
                         (r.font.name or "") .lower().startswith("consolas"))
                        for r in p.runs]
                if runs:
                    paras.append((p, runs))
            if not paras:
                continue
            pad = 4 if str(sh.shape_type).startswith("TEXT_BOX") else 8
            avail = (x1 - x0) - 2 * pad
            # layout
            blocks = []
            total_h = 0
            for p, runs in paras:
                base = max(r[1] for r in runs)
                ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.1
                lh = base * ls * 96 / 72 * scale
                joined = "".join(r[0] for r in runs)
                f0 = font(base, runs[0][3], runs[0][4], scale)
                lines = wrap(d, joined, f0, avail) if avail > 10 else [joined]
                sa = (p.space_after.pt if p.space_after else 3) * 96 / 72 * scale
                blocks.append((runs, lines, lh, sa, base))
                total_h += lh * len(lines) + sa
            va = tf.vertical_anchor
            cy = y0 + pad
            if va is not None and int(va) == 3:  # middle
                cy = y0 + ((y1 - y0) - total_h) / 2
            for runs, lines, lh, sa, base in blocks:
                align = paras[0][0].alignment
                for li, line in enumerate(lines):
                    # map runs onto this line greedily
                    seg, consumed = [], 0
                    remaining = line
                    pos = 0
                    # rebuild run segmentation across wrapped lines
                    seg = None
                    tx = x0 + pad
                    if li == 0 and len(runs) > 1 and len(lines) == 1:
                        for rt, rs, rc, rb, rm in runs:
                            f = font(rs, rb, rm, scale)
                            rt = rt.replace("\n", " ")
                            d.text((tx, cy), rt, font=f, fill=rc)
                            tx += d.textlength(rt, font=f)
                    else:
                        rt, rs, rc, rb, rm = runs[0]
                        f = font(base, rb, rm, scale)
                        w = d.textlength(line, font=f)
                        ax = x0 + pad
                        if align is not None and int(align) == 2:  # center
                            ax = x0 + ((x1 - x0) - w) / 2
                        elif align is not None and int(align) == 3:  # right
                            ax = x1 - pad - w
                        d.text((ax, cy), line, font=f, fill=rc)
                    cy += lh
                cy += sa
        out = os.path.join(outdir, f"slide_{idx:02d}.png")
        img.save(out)
        paths.append(out)
    return paths


if __name__ == "__main__":
    src = sys.argv[1]
    out = sys.argv[2]
    for p in render(src, out, float(sys.argv[3]) if len(sys.argv) > 3 else 1.0):
        print(p)
