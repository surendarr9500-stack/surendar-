"""Generate the CAPACITY CONNECT SIH 2026 (PS 26075) pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "CAPACITY_CONNECT_SIH2026_26075.pptx")

# ---------------- Theme ----------------
NAVY = RGBColor(0x06, 0x1B, 0x36)
NAVY_2 = RGBColor(0x0C, 0x2A, 0x4D)
TEAL = RGBColor(0x1FB, 0x00, 0x00) if False else RGBColor(0x21, 0xC7, 0xB8)
CYAN = RGBColor(0x3E, 0xC5, 0xF0)
SAFF = RGBColor(0xFF, 0x9A, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB9, 0xC9, 0xDA)
CARD = RGBColor(0x10, 0x2F, 0x53)
GREEN = RGBColor(0x4C, 0xD9, 0x7B)
RED = RGBColor(0xFF, 0x6B, 0x6B)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=CARD, line=None, radius=True, lw=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    first = True
    for item in runs:
        txt, size, color, bold = (list(item) + [False])[:4]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Segoe UI"
    return tb


def header(s, num, title, sub=None):
    bar = rect(s, 0, 0, W, Inches(1.05), NAVY_2, radius=False)
    rect(s, 0, Inches(1.05), W, Pt(3), TEAL, radius=False)
    text(s, Inches(0.55), Inches(0.2), Inches(0.8), Inches(0.6),
         [(num, 26, TEAL, True)])
    text(s, Inches(1.25), Inches(0.16), Inches(10.4), Inches(0.75),
         [(title, 26, WHITE, True)] + ([(sub, 12, MUTED, False)] if sub else []))
    text(s, Inches(11.0), Inches(0.32), Inches(1.9), Inches(0.4),
         [("SIH 2026 · PS 26075", 10, MUTED, False)], align=PP_ALIGN.RIGHT)
    return bar


def bullets(s, x, y, w, items, size=13, gap=0.42, bullet_color=TEAL):
    for i, it in enumerate(items):
        cy = y + Inches(gap * i)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, cy + Inches(0.055), Pt(7), Pt(7))
        d.fill.solid(); d.fill.fore_color.rgb = bullet_color
        d.line.fill.background(); d.shadow.inherit = False
        if isinstance(it, tuple):
            text(s, x + Inches(0.24), cy, w - Inches(0.24), Inches(0.4),
                 [(it[0] + "  ", size, WHITE, True), (it[1], size, MUTED, False)])
            tb = s.shapes[-1].text_frame.paragraphs[0]
        else:
            text(s, x + Inches(0.24), cy, w - Inches(0.24), Inches(0.4),
                 [(it, size, MUTED, False)])


def card(s, x, y, w, h, title, body, accent=TEAL, tsize=14, bsize=11):
    c = rect(s, x, y, w, h, CARD, line=RGBColor(0x1C, 0x44, 0x72))
    rect(s, x, y, Pt(4), h, accent, radius=False)
    text(s, x + Inches(0.22), y + Inches(0.16), w - Inches(0.4), Inches(0.4),
         [(title, tsize, accent, True)])
    text(s, x + Inches(0.22), y + Inches(0.16) + Pt(tsize + 10), w - Inches(0.4),
         h - Inches(0.5), [(b, bsize, MUTED, False) for b in body], space_after=3,
         line_spacing=1.15)
    return c


def chip(s, x, y, label, color=CYAN, w=None, size=10):
    w = w or Inches(0.12 * len(label) + 0.3)
    r = rect(s, x, y, w, Inches(0.32), None, line=color)
    text(s, x, y + Inches(0.04), w, Inches(0.25), [(label, size, color, True)],
         align=PP_ALIGN.CENTER)
    return x + w + Inches(0.12)


def arrow_row(s, x, y, w, labels, color=CYAN, size=10):
    n = len(labels)
    gap = Inches(0.18)
    bw = (w - gap * (n - 1)) / n
    for i, lab in enumerate(labels):
        bx = x + (bw + gap) * i
        b = rect(s, bx, y, bw, Inches(0.62), NAVY_2, line=color)
        text(s, bx + Inches(0.06), y + Inches(0.08), bw - Inches(0.12), Inches(0.46),
             [(lab, size, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, bx + bw + Inches(0.02),
                                   y + Inches(0.23), Inches(0.14), Inches(0.16))
            a.fill.solid(); a.fill.fore_color.rgb = TEAL
            a.line.fill.background(); a.shadow.inherit = False


# ======================= 1. TITLE =======================
s = slide()
band = rect(s, 0, 0, W, Inches(2.6), NAVY_2, radius=False)
rect(s, 0, Inches(2.6), W, Pt(4), TEAL, radius=False)
for i, (cx, sz, col) in enumerate([(11.4, 2.6, 0x14), (12.4, 1.7, 0x18), (10.6, 1.0, 0x16)]):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(3.4 - sz / 2),
                           Inches(sz), Inches(sz))
    o.fill.solid(); o.fill.fore_color.rgb = RGBColor(0x0C, 0x2A + i * 4, 0x4D)
    o.line.fill.background(); o.shadow.inherit = False

text(s, Inches(0.8), Inches(0.55), Inches(9), Inches(0.4),
     [("SMART INDIA HACKATHON 2026  ·  SOFTWARE  ·  SMART EDUCATION", 12, TEAL, True)])
text(s, Inches(0.8), Inches(1.05), Inches(11), Inches(1.2),
     [("CAPACITY CONNECT", 52, WHITE, True)])
text(s, Inches(0.8), Inches(1.95), Inches(11), Inches(0.5),
     [("Digital Capacity Building & Learning Management Portal", 19, CYAN, False)])

text(s, Inches(0.8), Inches(3.05), Inches(8.2), Inches(1.2),
     [("An offline-first, secure, role-based learning and competency platform for "
       "Ministry of Earth Sciences field personnel — training, assessments, trainer "
       "libraries, competency mapping and AI-assisted operational troubleshooting on "
       "one centralized portal.", 13.5, MUTED, False)], line_spacing=1.3)

x = Inches(0.8)
for lab, c in [("PS ID 26075", SAFF), ("Ministry of Earth Sciences", CYAN),
               ("India Meteorological Dept.", CYAN), ("Software", TEAL)]:
    x = Inches(chip(s, x, Inches(4.35), lab, c) / 914400)

for i, (k, v) in enumerate([("3", "User Roles"), ("100%", "Offline Capable"),
                            ("AES-256", "GCM Encryption"), ("4", "Target Platforms")]):
    cx = Inches(0.8 + i * 2.35)
    rect(s, cx, Inches(5.05), Inches(2.1), Inches(1.05), CARD,
         line=RGBColor(0x1C, 0x44, 0x72))
    text(s, cx, Inches(5.2), Inches(2.1), Inches(0.45), [(k, 22, TEAL, True)],
         align=PP_ALIGN.CENTER)
    text(s, cx, Inches(5.68), Inches(2.1), Inches(0.3), [(v, 10, MUTED, False)],
         align=PP_ALIGN.CENTER)

text(s, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.4),
     [("Team submission  ·  Flutter + FastAPI + PostgreSQL/SQLite  ·  Local AI Engine  ·  3D Digital Twin",
       11, MUTED, False)])

# ======================= 2. PROBLEM =======================
s = slide()
header(s, "01", "The Problem", "Why MoES / IMD capacity building needs a new platform")
text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5),
     [("Training, competency data and operational knowledge in earth-science "
       "organisations are fragmented across files, emails and disconnected portals — "
       "and unavailable exactly where field work happens.", 13, MUTED, False)],
     line_spacing=1.25)

probs = [
    ("Fragmented learning", ["Courses, lectures and study material scattered across "
                             "drives, mail and local machines.", "No single catalog, no versioning, no audit of who learned what."]),
    ("No competency map", ["Identifying a qualified trainer for a subject is manual and "
                           "opinion-driven.", "Skills, certificates and experience are not machine-readable."]),
    ("Connectivity blackspots", ["Research vessels, buoy sites and remote observatories "
                                 "operate with little or no Internet.", "Cloud-only portals become unusable at the point of need."]),
    ("Manual assessment", ["MCQ tests, deadlines and scorecards handled on paper or "
                           "spreadsheets.", "Slow feedback loops; participation statistics are guesswork."]),
    ("Weak knowledge reuse", ["Troubleshooting expertise for sonar, telemetry and Argo "
                              "assets lives in people's heads.", "New engineers repeat the same diagnostic learning curve."]),
    ("No governance", ["Approvals, role changes and content publishing lack a workflow "
                       "and audit trail.", "Admins have no live dashboard of enrollment or certification."]),
]
for i, (t, b) in enumerate(probs):
    col, row = i % 3, i // 3
    card(s, Inches(0.6 + col * 4.12), Inches(2.2 + row * 2.25), Inches(3.9),
         Inches(2.0), t, b, accent=SAFF if row == 0 else CYAN)

# ======================= 3. SOLUTION =======================
s = slide()
header(s, "02", "Proposed Solution", "One portal: learn, assess, certify, diagnose — online or offline")
text(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.4),
     [("CAPACITY CONNECT is a centralized, role-based capacity-building portal with an "
       "offline-first edge application. Everything a trainee, trainer or admin does is "
       "written to a real local database first, then synchronized to the cloud through an "
       "auditable transaction ledger.", 12.5, MUTED, False)], line_spacing=1.25)

roles = [
    ("TRAINEE", TEAL, ["Professional profile: qualifications, experience,",
                       "skills, interests, certificates",
                       "Course catalog, enrollment, resource access",
                       "Subject-wise MCQ assessments with scorecards",
                       "Course & content feedback",
                       "AI troubleshooting + 3D Digital Twin assist"]),
    ("TRAINER", CYAN, ["Trainer profile & competency declaration",
                       "Build questionnaires with deadlines",
                       "Monitor participation and performance",
                       "Trainer Library: recorded lectures, slides,",
                       "study material with versions & checksums",
                       "Per-course analytics and feedback review"]),
    ("ADMIN", SAFF, ["User approval and role management",
                     "Dashboards: courses, enrollments, certifications,",
                     "assessments, participation statistics",
                     "Publish notices, announcements, achievements",
                     "Competency mapping → best-fit trainer per subject",
                     "Assets, knowledge base, audit and sync control"]),
]
for i, (t, c, b) in enumerate(roles):
    card(s, Inches(0.6 + i * 4.12), Inches(2.35), Inches(3.9), Inches(2.75), t, b,
         accent=c, tsize=15, bsize=10.5)

rect(s, Inches(0.6), Inches(5.35), Inches(12.13), Inches(1.5), NAVY_2,
     line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(0.3),
     [("WHAT MAKES IT DIFFERENT", 11, TEAL, True)])
diffs = ["Offline-first by design — local SQLite/Drift is the operating DB, not a cache",
         "Local AI engine on 127.0.0.1 — deterministic retrieval works with zero cloud",
         "3D Digital Twin maps a described fault to a real mesh (SONAR-001 → Mesh_042)",
         "Competency mapping engine scores trainers against subject requirements"]
for i, d in enumerate(diffs):
    col, row = i % 2, i // 2
    text(s, Inches(0.9 + col * 5.9), Inches(5.9 + row * 0.38), Inches(5.7), Inches(0.35),
         [("▸  ", 11, TEAL, True), (d, 11, MUTED, False)])

# ======================= 4. ARCHITECTURE =======================
s = slide()
header(s, "03", "System Architecture", "Modular edge application + independently deployable cloud platform")

rect(s, Inches(0.6), Inches(1.5), Inches(5.85), Inches(4.5), NAVY_2,
     line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(0.85), Inches(1.65), Inches(5.3), Inches(0.3),
     [("EDGE APPLICATION  (works with no Internet)", 12, TEAL, True)])
edge = [("Flutter Client", "Material 3 · Riverpod · GoRouter · responsive mobile/desktop shells"),
        ("Local AI Engine", "Python FastAPI node on 127.0.0.1 · NLP pipeline · structured JSON"),
        ("Local Database", "Drift/SQLite (SQLCipher) · users, courses, attempts, diagnostics, queue"),
        ("Digital Twin Runtime", "Locally cached GLTF/GLB · selection, highlight, fault state"),
        ("Secure Storage", "Platform keystore · AES-256-GCM data protection · session tokens")]
for i, (t, d) in enumerate(edge):
    y = Inches(2.1 + i * 0.76)
    rect(s, Inches(0.85), y, Inches(5.35), Inches(0.66), CARD, line=None)
    text(s, Inches(1.0), y + Inches(0.06), Inches(5.1), Inches(0.28), [(t, 11.5, WHITE, True)])
    text(s, Inches(1.0), y + Inches(0.33), Inches(5.1), Inches(0.28), [(d, 9.5, MUTED, False)])

rect(s, Inches(6.9), Inches(1.5), Inches(5.83), Inches(4.5), NAVY_2,
     line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(7.15), Inches(1.65), Inches(5.3), Inches(0.3),
     [("CLOUD PLATFORM  (optional for core workflow)", 12, SAFF, True)])
cloud = [("REST API", "FastAPI · versioned /api/v1 · Pydantic validation · consistent errors"),
         ("PostgreSQL", "SQLAlchemy models · Alembic migrations · full relational schema"),
         ("Background Workers", "Media transcode, document indexing, report generation"),
         ("Sync & Conflict Service", "Ledger ingestion · version compare · conflict resolution"),
         ("Admin Portal", "Approvals, dashboards, publishing, audit log, device registry")]
for i, (t, d) in enumerate(cloud):
    y = Inches(2.1 + i * 0.76)
    rect(s, Inches(7.15), y, Inches(5.35), Inches(0.66), CARD, line=None)
    text(s, Inches(7.3), y + Inches(0.06), Inches(5.1), Inches(0.28), [(t, 11.5, WHITE, True)])
    text(s, Inches(7.3), y + Inches(0.33), Inches(5.1), Inches(0.28), [(d, 9.5, MUTED, False)])

sy = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(6.2),
                        Inches(4.1), Inches(0.62))
sy.fill.solid(); sy.fill.fore_color.rgb = TEAL; sy.line.fill.background()
sy.shadow.inherit = False
text(s, Inches(4.6), Inches(6.35), Inches(4.1), Inches(0.35),
     [("⇅  ENCRYPTED SYNC LEDGER  ⇅", 13, NAVY, True)], align=PP_ALIGN.CENTER)

# ======================= 5. TECH STACK =======================
s = slide()
header(s, "04", "Technology Stack", "Chosen for cross-platform reach, offline strength and low deployment cost")
groups = [
    ("CLIENT", TEAL, ["Flutter 3 / Dart", "Material 3 design system", "Riverpod state management",
                      "GoRouter navigation", "Drift + SQLite (SQLCipher)",
                      "flutter_secure_storage", "Dio HTTP client", "connectivity_plus"]),
    ("3D & MEDIA", CYAN, ["GLTF / GLB assets", "flutter_scene / model viewer", "Local asset cache + checksum",
                          "Chunked, resumable downloads", "Streaming video playback",
                          "PDF / document viewer", "speech_to_text (optional voice)"]),
    ("AI ENGINE", SAFF, ["Python 3.11 · FastAPI", "Normalization + language detect", "Tokenization & stop-wording",
                         "Keyword / phrase / fuzzy match", "RapidFuzz scoring", "TF-IDF knowledge retrieval",
                         "Pluggable local LLM slot"]),
    ("BACKEND & OPS", GREEN, ["FastAPI · SQLAlchemy · Pydantic", "PostgreSQL + Alembic", "Celery / RQ workers",
                              "JWT access + refresh tokens", "Docker + docker-compose", "GitHub Actions CI",
                              "pytest · flutter test"]),
]
for i, (t, c, items) in enumerate(groups):
    x = Inches(0.6 + i * 3.11)
    rect(s, x, Inches(1.5), Inches(2.95), Inches(4.5), CARD, line=RGBColor(0x1C, 0x44, 0x72))
    rect(s, x, Inches(1.5), Inches(2.95), Pt(4), c, radius=False)
    text(s, x + Inches(0.2), Inches(1.72), Inches(2.6), Inches(0.3), [(t, 13, c, True)])
    for j, it in enumerate(items):
        text(s, x + Inches(0.2), Inches(2.2 + j * 0.43), Inches(2.6), Inches(0.35),
             [("· ", 11, c, True), (it, 10.5, MUTED, False)], line_spacing=1.1)

rect(s, Inches(0.6), Inches(6.2), Inches(12.13), Inches(0.75), NAVY_2,
     line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(0.85), Inches(6.32), Inches(11.6), Inches(0.5),
     [("Targets: ", 11.5, TEAL, True),
      ("Android (field tablets/phones) · Windows (vessel workstations) · Linux (shore servers/kiosks) · Web (admin & trainer portal). "
       "One Dart codebase, adaptive layouts, single local schema.", 11, MUTED, False)],
     line_spacing=1.2)

# ======================= 6. AI ENGINE =======================
s = slide()
header(s, "05", "Local AI Troubleshooting Engine", "Runs on 127.0.0.1 — no cloud LLM required for the core workflow")
arrow_row(s, Inches(0.6), Inches(1.45), Inches(12.13),
          ["Normalize", "Language\nDetect", "Tokenize", "Keyword\nMatch", "Phrase\nMatch",
           "Fuzzy\nMatch", "Knowledge\nRetrieval"], size=9)
arrow_row(s, Inches(2.4), Inches(2.2), Inches(8.5),
          ["Component\nID", "Fault\nClassify", "Severity", "Recommended\nAction", "3D Mesh\nMapping"], size=9)

card(s, Inches(0.6), Inches(3.15), Inches(5.85), Inches(1.55), "SAMPLE INPUT",
     ['"Sonar transducer is showing abnormal vibration and casing fracture."',
      "",
      "Sonar → Sonar Transducer Array → SONAR-001 → Mesh_042 → HIGH severity → diagnostic procedure"],
     accent=SAFF, bsize=10.5)

rect(s, Inches(6.9), Inches(3.15), Inches(5.83), Inches(3.55), RGBColor(0x08, 0x21, 0x3E),
     line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(7.15), Inches(3.3), Inches(5.3), Inches(0.3),
     [("STRUCTURED JSON RESPONSE", 11, TEAL, True)])
code = ['{', '  "request_id": "b7c1…",', '  "component_id": "SONAR-001",',
        '  "component_name": "Sonar Transducer Array",', '  "mesh_id": "Mesh_042",',
        '  "fault": "Casing fracture",', '  "severity": "HIGH",', '  "confidence": 0.94,',
        '  "evidence": ["casing fracture", "vibration"],',
        '  "recommended_actions": ["Isolate array", …],',
        '  "warnings": ["Do not power-cycle underwater"],', '  "timestamp": "2026-09-01T10:22:04Z"', '}']
tb = s.shapes.add_textbox(Inches(7.15), Inches(3.65), Inches(5.4), Inches(2.9))
tf = tb.text_frame; tf.word_wrap = True
for i, ln in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1)
    r = p.add_run(); r.text = ln
    r.font.size = Pt(10); r.font.name = "Consolas"
    r.font.color.rgb = CYAN if '"' in ln else MUTED

card(s, Inches(0.6), Inches(4.9), Inches(5.85), Inches(1.8), "HOW CONFIDENCE IS COMPUTED",
     ["Weighted score, not a magic number:",
      "0.45 × exact keyword coverage  +  0.25 × phrase match  +  0.20 × fuzzy token ratio  +  0.10 × knowledge-chunk cosine similarity",
      "Score < 0.55 → engine returns candidates and asks the engineer to disambiguate."],
     accent=GREEN, bsize=10.5)

# ======================= 7. DIGITAL TWIN =======================
s = slide()
header(s, "06", "Digital Twin & Hardware Registry", "Fault text → real component → real mesh → visual fault state")
arrow_row(s, Inches(0.6), Inches(1.45), Inches(12.13),
          ["Component\nRegistry", "Mesh\nMapping", "3D Scene\n(GLTF/GLB)", "Component\nSelection",
           "State Layer", "Visual\nHighlight"], size=9.5)

hdr = ["Component ID", "Name", "Mesh ID", "Category", "Demo State"]
rows = [["SONAR-001", "Sonar Transducer Array", "Mesh_042", "Acoustics", "CRITICAL"],
        ["TELEM-001", "Telemetry Transceiver Mast", "Mesh_109", "Communications", "WARNING"],
        ["ARGO-001", "Autonomous Argo Profiling Float", "Mesh_210", "Profiling", "NORMAL"],
        ["ECHO-001", "Multi-beam Echo Sounder", "Mesh_315", "Bathymetry", "DEGRADED"],
        ["WINCH-001", "Hydraulic Deep-Sea Winch", "Mesh_410", "Deck Machinery", "MAINTENANCE"]]
colw = [Inches(1.7), Inches(3.3), Inches(1.4), Inches(1.7), Inches(1.6)]
x0, y0 = Inches(0.6), Inches(2.35)
rect(s, x0, y0, Inches(9.7), Inches(0.42), NAVY_2, radius=False)
cx = x0
for i, hcell in enumerate(hdr):
    text(s, cx + Inches(0.12), y0 + Inches(0.1), colw[i], Inches(0.3), [(hcell, 10.5, TEAL, True)])
    cx += colw[i]
statecol = {"CRITICAL": RED, "WARNING": SAFF, "NORMAL": GREEN, "DEGRADED": SAFF,
            "MAINTENANCE": CYAN}
for r, row in enumerate(rows):
    ry = y0 + Inches(0.42 + r * 0.44)
    rect(s, x0, ry, Inches(9.7), Inches(0.42),
         CARD if r % 2 == 0 else RGBColor(0x0B, 0x26, 0x45), radius=False)
    cx = x0
    for i, cell in enumerate(row):
        col = statecol.get(cell, WHITE if i == 0 else MUTED)
        text(s, cx + Inches(0.12), ry + Inches(0.11), colw[i], Inches(0.3),
             [(cell, 10.5, col, i in (0, 4))])
        cx += colw[i]

card(s, Inches(10.5), Inches(2.35), Inches(2.23), Inches(2.62), "STATES",
     ["NORMAL", "WARNING", "DEGRADED", "CRITICAL", "MAINTENANCE", "OFFLINE", "UNKNOWN"],
     accent=CYAN, bsize=10.5)

feat = [("Interaction", "rotate · zoom · pan · select · isolate · reset camera"),
        ("Offline assets", "GLTF/GLB cached locally with checksum + version"),
        ("Decoupled state", "AI parser never touches the renderer — a state layer sits between"),
        ("Graceful fallback", "model missing → full component metadata still accessible")]
for i, (t, d) in enumerate(feat):
    x = Inches(0.6 + (i % 2) * 6.11)
    y = Inches(5.05 + (i // 2) * 0.9)
    rect(s, x, y, Inches(5.9), Inches(0.78), CARD, line=None)
    text(s, x + Inches(0.18), y + Inches(0.1), Inches(5.5), Inches(0.28), [(t, 11.5, TEAL, True)])
    text(s, x + Inches(0.18), y + Inches(0.4), Inches(5.5), Inches(0.28), [(d, 10, MUTED, False)])

# ======================= 8. OFFLINE & SYNC =======================
s = slide()
header(s, "07", "Offline-First Operation & Synchronization", "The local database is an operating database, not a cache")
arrow_row(s, Inches(0.6), Inches(1.45), Inches(12.13),
          ["Local\nTransaction", "Sync Queue", "Connectivity\nRestored", "Re-auth",
           "Upload", "Server\nValidation", "Conflict\nDetection", "ACK →\nSYNCED"], size=9)

card(s, Inches(0.6), Inches(2.35), Inches(3.9), Inches(2.3), "WORKS FULLY OFFLINE",
     ["Policy-based offline authentication", "Downloaded courses, lectures, documents",
      "MCQ assessments and scoring", "AI troubleshooting + Digital Twin",
      "Diagnostic records with attachments", "Local knowledge search"], accent=TEAL)

card(s, Inches(4.72), Inches(2.35), Inches(3.9), Inches(2.3), "LEDGER RECORD",
     ["transaction_id · device_id · user_id", "entity_type · entity_id · operation",
      "payload · created_at · updated_at", "sync_status · retry_count",
      "States: PENDING → SYNCING → SYNCED", "                    ↘ FAILED / CONFLICT"],
     accent=CYAN)

card(s, Inches(8.83), Inches(2.35), Inches(3.9), Inches(2.3), "CONFLICT RESOLUTION",
     ["Per-entity strategy — never blind overwrite:", "Progress / quiz → highest-attempt merge",
      "Profile fields → field-level merge", "Course content → server version wins",
      "Diagnostics → manual resolution queue", "Critical records are never dropped silently"],
     accent=SAFF)

rect(s, Inches(0.6), Inches(4.85), Inches(12.13), Inches(1.9), NAVY_2,
     line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.3),
     [("FAILURE BEHAVIOUR — tested explicitly", 11.5, TEAL, True)])
fails = [("Internet drops mid-session", "banner flips to OFFLINE · LOCAL ENGINE ACTIVE"),
         ("Backend crashes", "edge app continues; queue keeps growing safely"),
         ("Local AI unavailable", "deterministic rule-based fallback path"),
         ("App killed during sync", "idempotent transaction IDs, queue replays"),
         ("Corrupted asset", "checksum mismatch → re-download, old version retained"),
         ("Invalid API response", "typed error envelope → user-visible validation message")]
for i, (a, b) in enumerate(fails):
    col, row = i % 3, i // 3
    x = Inches(0.9 + col * 4.0)
    y = Inches(5.45 + row * 0.55)
    text(s, x, y, Inches(3.8), Inches(0.5),
         [(a, 10.5, WHITE, True), (b, 9.5, MUTED, False)], line_spacing=1.15)

# ======================= 9. SECURITY =======================
s = slide()
header(s, "08", "Security & Governance", "Defence in layers — enforced on the server, not just hidden in the UI")
layers = [("LOCAL", TEAL, ["Encrypted database (SQLCipher)",
                           "Keys in platform keystore / DPAPI",
                           "AES-256-GCM for sensitive blobs",
                           "Session timeout & re-auth policy",
                           "Integrity checks on assets"]),
          ("NETWORK", CYAN, ["HTTPS with certificate validation",
                             "Authenticated, versioned API",
                             "Request schema validation",
                             "Rate limiting per device/user",
                             "Device registration & revocation"]),
          ("BACKEND", SAFF, ["RBAC on every endpoint",
                             "Argon2id password hashing",
                             "JWT access + rotating refresh",
                             "Input validation via Pydantic",
                             "No secrets in code — .env only"]),
          ("AUDIT", GREEN, ["LOGIN / LOGOUT",
                            "DOCUMENT_ACCESS · AI_ANALYSIS",
                            "TRAINING_COMPLETED · QUIZ_ATTEMPT",
                            "DIAGNOSTIC_CREATED / UPDATED",
                            "SYNC_* · ADMIN_ACTION"])]
for i, (t, c, items) in enumerate(layers):
    x = Inches(0.6 + i * 3.11)
    card(s, x, Inches(1.5), Inches(2.95), Inches(2.6), t, items, accent=c, tsize=13)

rect(s, Inches(0.6), Inches(4.3), Inches(5.85), Inches(2.45), CARD, line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(0.85), Inches(4.45), Inches(5.4), Inches(0.3),
     [("ROLE-BASED PERMISSIONS", 12, TEAL, True)])
perm = [("Trainee / Field Engineer", "profile · enroll · learn · assess · feedback · diagnostics"),
        ("Trainer / Training Officer", "content · questionnaires · deadlines · library · analytics"),
        ("Admin", "approvals · roles · dashboards · publishing · assets · security"),
        ("Supervisor (extension)", "read-only oversight of team progress and asset health")]
for i, (r, d) in enumerate(perm):
    y = Inches(4.9 + i * 0.45)
    text(s, Inches(0.85), y, Inches(5.4), Inches(0.4),
         [(r + " → ", 10.5, WHITE, True), (d, 10, MUTED, False)], line_spacing=1.1)

rect(s, Inches(6.9), Inches(4.3), Inches(5.83), Inches(2.45), CARD, line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(7.15), Inches(4.45), Inches(5.4), Inches(0.3),
     [("OFFLINE AUTHENTICATION POLICY", 12, SAFF, True)])
text(s, Inches(7.15), Inches(4.9), Inches(5.35), Inches(1.7),
     [("Offline login is a deliberate, bounded compromise and is documented as such:", 10.5, MUTED, False),
      ("· Only devices registered and approved by an Admin may cache credentials.", 10.5, MUTED, False),
      ("· A salted Argon2id verifier — never the password — is cached locally.", 10.5, MUTED, False),
      ("· The offline grace window is configurable (default 7 days) then hard-locks.", 10.5, MUTED, False),
      ("· All offline actions are audit-logged and re-verified at next sync.", 10.5, MUTED, False)],
     space_after=4, line_spacing=1.15)

# ======================= 10. DATA MODEL =======================
s = slide()
header(s, "09", "Data Model & API Surface", "Migration-driven schema, versioned REST endpoints")
domains = [("IDENTITY", ["users", "roles", "permissions", "devices", "sessions", "profiles",
                         "qualifications", "certificates"]),
           ("LEARNING", ["courses", "modules", "lessons", "media", "documents", "enrollments",
                         "progress", "feedback"]),
           ("ASSESSMENT", ["quizzes", "questions", "options", "attempts", "answers", "scores",
                           "deadlines", "certifications"]),
           ("OPERATIONS", ["components", "component_faults", "maintenance_procedures",
                           "digital_twin_models", "diagnostics", "maintenance_records",
                           "work_orders", "attachments"]),
           ("PLATFORM", ["sync_transactions", "audit_logs", "notifications", "announcements",
                         "competency_map", "trainer_subjects", "settings", "content_versions"])]
for i, (t, tables) in enumerate(domains):
    x = Inches(0.6 + i * 2.46)
    rect(s, x, Inches(1.5), Inches(2.3), Inches(3.15), CARD, line=RGBColor(0x1C, 0x44, 0x72))
    rect(s, x, Inches(1.5), Inches(2.3), Pt(4), TEAL if i % 2 == 0 else CYAN, radius=False)
    text(s, x + Inches(0.16), Inches(1.7), Inches(2.0), Inches(0.3), [(t, 11.5, WHITE, True)])
    for j, tb_ in enumerate(tables):
        text(s, x + Inches(0.16), Inches(2.12 + j * 0.31), Inches(2.05), Inches(0.28),
             [(tb_, 9.5, MUTED, False)])

rect(s, Inches(0.6), Inches(4.85), Inches(12.13), Inches(1.9), NAVY_2, line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.3),
     [("VERSIONED API  /api/v1", 11.5, SAFF, True)])
eps = ["/auth", "/users", "/profiles", "/courses", "/lessons", "/media", "/documents",
       "/quizzes", "/attempts", "/enrollments", "/components", "/digital-twin",
       "/diagnostics", "/maintenance", "/competency", "/notifications", "/sync", "/audit"]
x = Inches(0.9); y = Inches(5.45)
for e in eps:
    w = Inches(0.085 * len(e) + 0.4)
    if x + w > Inches(12.5):
        x = Inches(0.9); y = y + Inches(0.5)
    rect(s, x, y, w, Inches(0.36), CARD, line=RGBColor(0x2A, 0x5A, 0x8C))
    text(s, x, y + Inches(0.07), w, Inches(0.26), [(e, 10, CYAN, False)], align=PP_ALIGN.CENTER)
    x = x + w + Inches(0.12)

# ======================= 11. DEMO FLOW =======================
s = slide()
header(s, "10", "Live Engineering Demo", "The demo is the real product pipeline — no scripted animation")
steps = [("1", "Login", "Trainee credentials, session issued"),
         ("2", "Dashboard", "Live counts from the local DB"),
         ("3", "Disable Internet", "UI flips to OFFLINE · LOCAL ENGINE ACTIVE"),
         ("4", "Enter Fault", '"abnormal vibration and casing fracture"'),
         ("5", "Local AI", "127.0.0.1 pipeline returns structured JSON"),
         ("6", "Component ID", "SONAR-001 · confidence 0.94 · HIGH"),
         ("7", "Digital Twin", "Mesh_042 highlighted in CRITICAL state"),
         ("8", "Guidance", "Diagnostic procedure from knowledge base"),
         ("9", "Create Record", "Diagnostic saved locally with attachments"),
         ("10", "Take Quiz", "MCQ attempt scored and persisted offline"),
         ("11", "Restart App", "All data still present — proves real persistence"),
         ("12", "Reconnect & Sync", "Queue drains, records marked SYNCED, admin sees them")]
for i, (n, t, d) in enumerate(steps):
    col, row = i % 4, i // 4
    x = Inches(0.6 + col * 3.11)
    y = Inches(1.55 + row * 1.42)
    rect(s, x, y, Inches(2.95), Inches(1.22), CARD, line=RGBColor(0x1C, 0x44, 0x72))
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15),
                              Inches(0.4), Inches(0.4))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background(); circ.shadow.inherit = False
    text(s, x + Inches(0.15), y + Inches(0.22), Inches(0.4), Inches(0.3),
         [(n, 12, NAVY, True)], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.68), y + Inches(0.2), Inches(2.1), Inches(0.3), [(t, 12, WHITE, True)])
    text(s, x + Inches(0.18), y + Inches(0.63), Inches(2.6), Inches(0.5),
         [(d, 9.5, MUTED, False)], line_spacing=1.15)

rect(s, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.72), NAVY_2, line=TEAL)
text(s, Inches(0.85), Inches(6.2), Inches(11.6), Inches(0.4),
     [("Acceptance criterion: ", 11.5, TEAL, True),
      ("every stage above is executed against the real database, the real parser and the real sync ledger. If any critical stage fails, the build is not shipped.",
       11, MUTED, False)])

# ======================= 12. FEASIBILITY =======================
s = slide()
header(s, "11", "Feasibility, Risks & Mitigation", "Honest engineering assessment")
rows = [("Large 3D models on low-end field tablets",
         "Decimated LOD meshes shipped by default; full-resolution GLB is an optional download; renderer degrades to metadata-only view."),
        ("No local LLM available on device",
         "Deterministic keyword/phrase/fuzzy + TF-IDF retrieval is the baseline and is fully functional; LLM slot is pluggable, not required."),
        ("Offline credential caching weakens auth",
         "Admin-approved devices only, Argon2id verifier, bounded grace window, full audit and re-verification on reconnect."),
        ("Sync conflicts on shared records",
         "Per-entity strategies with an explicit manual-resolution queue; diagnostics are never auto-overwritten."),
        ("Storage exhaustion on field devices",
         "Storage manager with per-category usage, resumable downloads, permissioned removal of optional assets only."),
        ("Scope is large for hackathon timeline",
         "Strict priority order: working app → offline → troubleshooting → twin → diagnostics → training → sync → security → admin.")]
y = Inches(1.5)
rect(s, Inches(0.6), y, Inches(12.13), Inches(0.42), NAVY_2, radius=False)
text(s, Inches(0.8), y + Inches(0.1), Inches(4.4), Inches(0.3), [("RISK / CHALLENGE", 11, TEAL, True)])
text(s, Inches(5.4), y + Inches(0.1), Inches(7.0), Inches(0.3), [("MITIGATION", 11, TEAL, True)])
for i, (r, m) in enumerate(rows):
    ry = y + Inches(0.42 + i * 0.82)
    rect(s, Inches(0.6), ry, Inches(12.13), Inches(0.8),
         CARD if i % 2 == 0 else RGBColor(0x0B, 0x26, 0x45), radius=False)
    text(s, Inches(0.8), ry + Inches(0.16), Inches(4.4), Inches(0.6),
         [(r, 10.5, WHITE, True)], line_spacing=1.15)
    text(s, Inches(5.4), ry + Inches(0.16), Inches(7.1), Inches(0.6),
         [(m, 10, MUTED, False)], line_spacing=1.15)

text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
     [("Viability: ", 11, TEAL, True),
      ("entirely open-source stack, single Dart codebase across four targets, containerized backend deployable on existing MoES/IMD infrastructure — no per-seat licensing.",
       10.5, MUTED, False)])

# ======================= 13. IMPACT =======================
s = slide()
header(s, "12", "Impact & Benefits", "For personnel, for the organisation, for the mission")
groups = [("FOR TRAINEES / FIELD STAFF", TEAL,
           ["Learn and get certified from a vessel or remote observatory with zero connectivity",
            "Instant AI-assisted diagnostic guidance instead of waiting for shore support",
            "A portable, verifiable competency profile with skills and certificates"]),
          ("FOR TRAINERS", CYAN,
           ["One library for lectures, slides and material with versioning and reach analytics",
            "Questionnaires with deadlines, auto-scoring and participation visibility",
            "Feedback loop that shows which content actually works"]),
          ("FOR ADMINISTRATION", SAFF,
           ["Live dashboards of enrollment, certification and assessment across the organisation",
            "Competency mapping surfaces the right trainer for each subject, objectively",
            "Complete audit trail for governance, compliance and reporting"]),
          ("FOR THE MISSION", GREEN,
           ["Reduced instrument downtime through faster, guided fault resolution",
            "Institutional knowledge captured once and reused everywhere",
            "Scalable to other MoES bodies — INCOIS, NCPOR, NIOT — with the same core"])]
for i, (t, c, items) in enumerate(groups):
    col, row = i % 2, i // 2
    card(s, Inches(0.6 + col * 6.11), Inches(1.5 + row * 2.3), Inches(5.9), Inches(2.1),
         t, items, accent=c, tsize=14, bsize=11)

rect(s, Inches(0.6), Inches(6.15), Inches(12.13), Inches(0.85), NAVY_2, line=TEAL)
metrics = [("↓ 40%", "target fault-resolution time"), ("↑ 3×", "assessment throughput"),
           ("0", "cloud dependency offline"), ("1", "source of truth for competency")]
for i, (k, v) in enumerate(metrics):
    x = Inches(0.9 + i * 3.05)
    text(s, x, Inches(6.28), Inches(2.9), Inches(0.3), [(k, 15, TEAL, True)], align=PP_ALIGN.CENTER)
    text(s, x, Inches(6.62), Inches(2.9), Inches(0.3), [(v, 9.5, MUTED, False)], align=PP_ALIGN.CENTER)

# ======================= 14. ROADMAP =======================
s = slide()
header(s, "13", "Delivery Roadmap", "Phase-gated: plan → build → run → test → verify → next")
phases = [("P1–P4", "Foundation", "Repo, environment, Flutter shell, design system, routing", TEAL),
          ("P5–P8", "Core Platform", "Local database, authentication, RBAC, live dashboard", CYAN),
          ("P9–P12", "Learning", "Training, media management, documents, quiz engine", SAFF),
          ("P13–P18", "Intelligence", "Diagnostics, hardware registry, local AI, knowledge engine, Digital Twin integration", GREEN),
          ("P19–P22", "Connectivity", "Offline engine, synchronization, backend, admin portal", CYAN),
          ("P23–P32", "Hardening", "Security, voice, localization, reporting, observability, tests, production builds", SAFF)]
y = Inches(1.55)
for i, (p, t, d, c) in enumerate(phases):
    ry = y + Inches(i * 0.85)
    rect(s, Inches(0.6), ry, Inches(12.13), Inches(0.72), CARD, line=None)
    rect(s, Inches(0.6), ry, Pt(4), Inches(0.72), c, radius=False)
    text(s, Inches(0.85), ry + Inches(0.2), Inches(1.3), Inches(0.35), [(p, 13, c, True)])
    text(s, Inches(2.3), ry + Inches(0.2), Inches(2.6), Inches(0.35), [(t, 12.5, WHITE, True)])
    text(s, Inches(5.0), ry + Inches(0.22), Inches(7.5), Inches(0.35), [(d, 10.5, MUTED, False)])

text(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.4),
     [("Definition of Done: ", 11, TEAL, True),
      ("code exists AND build succeeds AND tests pass AND UI works AND database works AND error states work AND offline behaviour is verified AND documentation is updated.",
       10.5, MUTED, False)])

# ======================= 15. CLOSING =======================
s = slide()
rect(s, 0, 0, W, Inches(2.2), NAVY_2, radius=False)
rect(s, 0, Inches(2.2), W, Pt(4), TEAL, radius=False)
text(s, Inches(0.8), Inches(0.7), Inches(11), Inches(0.9), [("CAPACITY CONNECT", 44, WHITE, True)])
text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
     [("Build the capability, not just the portal.", 17, TEAL, False)])

text(s, Inches(0.8), Inches(2.8), Inches(7.3), Inches(1.6),
     [("A centralized, secure, role-based capacity-building platform that keeps working "
       "where MoES and IMD personnel actually work — on vessels, at buoy stations and in "
       "remote observatories — and reconciles everything the moment the network returns.",
       14, MUTED, False)], line_spacing=1.35)

items = ["Real offline-first local database, not a cache",
         "Local AI engine with a defined confidence algorithm",
         "Digital Twin driven by a decoupled state layer",
         "Auditable sync ledger with explicit conflict handling",
         "RBAC enforced at the API, not just in the UI"]
for i, it in enumerate(items):
    text(s, Inches(0.85), Inches(4.55 + i * 0.42), Inches(7.2), Inches(0.35),
         [("✓  ", 12, GREEN, True), (it, 12, MUTED, False)])

rect(s, Inches(8.6), Inches(2.8), Inches(4.13), Inches(3.9), CARD, line=RGBColor(0x1C, 0x44, 0x72))
text(s, Inches(8.9), Inches(3.05), Inches(3.6), Inches(0.3), [("SUBMISSION", 11, TEAL, True)])
meta = [("Problem Statement ID", "26075"),
        ("Title", "Capacity Connect — Digital Capacity Building & Learning Management Portal"),
        ("Organization", "Ministry of Earth Sciences"),
        ("Department", "India Meteorological Department"),
        ("Category", "Software"),
        ("Theme", "Smart Education")]
yy = Inches(3.5)
for k, v in meta:
    text(s, Inches(8.9), yy, Inches(3.6), Inches(0.28), [(k.upper(), 8.5, TEAL, True)])
    tb = text(s, Inches(8.9), yy + Inches(0.22), Inches(3.6), Inches(0.4), [(v, 11, WHITE, False)],
              line_spacing=1.1)
    yy = yy + Inches(0.52 if len(v) < 40 else 0.72)

prs.save(OUT)
print("saved", OUT, len(prs.slides.__iter__.__self__._sldIdLst))
