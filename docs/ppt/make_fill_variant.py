"""Generate `fill_official_template.py` from `build_sih_official.py`.

The two scripts share all layout/content code. This transform rewrites the
standalone builder into a version that opens the user's own SIH template
`.pptx`, clears the placeholder prompts and writes the same content into it,
scaling coordinates to whatever slide size that template uses.

Run after editing build_sih_official.py:

    python3 make_fill_variant.py
"""
import os
import re

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "build_sih_official.py")
DST = os.path.join(HERE, "fill_official_template.py")

s = open(SRC).read()

# ---------------- 1. header ----------------
NEW_HEAD = '''"""Fill the OFFICIAL SIH template (your own .pptx) with the Capacity Connect content.

This edits *your* template file: it keeps the file's masters, layouts, theme,
SIH logo, footer and slide titles, deletes only the placeholder prompt text,
and writes the Capacity Connect content and images into the body of each slide.

GENERATED FILE - edit build_sih_official.py and re-run make_fill_variant.py.

Usage
-----
    pip install python-pptx pillow
    python3 fill_official_template.py <your_template.pptx> [output.pptx]
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches as _Inches, Pt as _Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "sih_template.pptx")
OUT = sys.argv[2] if len(sys.argv) > 2 else os.path.join(
    HERE, "CAPACITY_CONNECT_26075_FILLED.pptx")

if not os.path.exists(TEMPLATE):
    sys.exit("Template not found: %s\\n"
             "Download your SIH template from Google Drive as .pptx "
             "(File -> Download -> Microsoft PowerPoint) and pass its path."
             % TEMPLATE)

prs = Presentation(TEMPLATE)

# Layout below is authored on a 13.333 x 7.5 in canvas. Google Slides exports
# 10 x 5.625 in - same 16:9 ratio, so scale everything uniformly.
DESIGN_W = 13.333
SCALE = (prs.slide_width / 914400.0) / DESIGN_W


def Inches(v):
    return Emu(int(round(v * SCALE * 914400)))


def Pt(v):
    return _Pt(v * SCALE)


def FS(v):
    """Font size in points, scaled to the template's canvas."""
    return _Pt(round(v * SCALE, 1))


W, H = prs.slide_width, prs.slide_height

# Placeholder text used by the blank SIH template - cleared before writing.
PROMPTS = [
    "problem statement id", "problem statement title", "theme-", "ps category",
    "team id", "team name (registered on portal)",
    "proposed solution", "detailed explanation of the proposed solution",
    "how it addresses the problem", "innovation and uniqueness",
    "technologies to be used", "methodology and process for implementation",
    "analysis of the feasibility", "potential challenges and risks",
    "strategies for overcoming", "potential impact on the target audience",
    "benefits of the solution", "details / links of the reference",
    "your team name",
]


def clear_body(slide):
    """Delete the template's prompt text, keep chrome.

    Shapes in the top 13% / bottom 13% of the slide (logo, team box, title,
    footer, page number) are never touched.
    """
    top_zone = H * 0.13
    bot_zone = H * 0.87
    for sh in list(slide.shapes):
        if not sh.has_text_frame or sh.top is None:
            continue
        if sh.top < top_zone or sh.top > bot_zone:
            continue
        txt = sh.text_frame.text.strip().lower()
        if not txt or any(k in txt for k in PROMPTS):
            sh._element.getparent().remove(sh._element)


'''
s = NEW_HEAD + s[s.index("# ---- SIH template palette"):]

# ---------------- 2. drop the standalone Presentation setup ----------------
s = re.sub(r"W, H = Inches\(13\.333\), Inches\(7\.5\)\nprs = Presentation\(\)\n"
           r"prs\.slide_width, prs\.slide_height = W, H\nBLANK = prs\.slide_layouts\[6\]\n\n",
           "", s)

# ---------------- 3. scale font sizes ----------------
s = s.replace("            r.font.size = Pt(size)", "            r.font.size = FS(size)")
s = s.replace("    p.space_after = Pt(space_after)", "    p.space_after = FS(space_after)")
s = s.replace("        p.space_after = Pt(space_after)", "        p.space_after = FS(space_after)")
s = s.replace('        r.font.size = Pt(8); r.font.color.rgb = GREY; r.font.name = "Calibri"',
              '        r.font.size = FS(8); r.font.color.rgb = GREY; r.font.name = "Calibri"')

# ---------------- 4. chrome()/slide() -> reuse the template's own slides ----
a = s.index("def chrome(s, num, title=None):")
b = s.index("# --- text measurement")
s = s[:a] + '''def slide(title=None, num=None):
    """Return the template's own slide `num`, cleared of prompt text.

    The template's chrome (SIH logo, footer, page number, team box, slide
    title) is left exactly as it is in your file.
    """
    idx = num - 1
    if idx >= len(prs.slides._sldIdLst):
        raise SystemExit("Your template has only %d slides; expected at least 6."
                         % len(prs.slides._sldIdLst))
    s = prs.slides[idx]
    clear_body(s)
    return s


''' + s[b:]

# ---------------- 5. title slide: keep the template's own headings ----------
a = s.index("# ==================== SLIDE 1 — TITLE PAGE ====================")
b = s.index("# ==================== SLIDE 2")
s = s[:a] + '''# ==================== SLIDE 1 — TITLE PAGE ====================
s = prs.slides[0]
clear_body(s)

fields = [("Problem Statement ID", "26075"),
          ("Problem Statement Title",
           "CAPACITY CONNECT — A Digital Capacity Building and Learning Management Portal"),
          ("Theme", "Smart Education"),
          ("PS Category", "Software"),
          ("Team ID", TEAM_ID),
          ("Team Name (Registered on portal)", TEAM)]
y = Inches(2.30)
for i, (k, v) in enumerate(fields):
    h = Inches(0.56)
    rect(s, Inches(0.7), y, Inches(8.2), h, LTGREY if i % 2 == 0 else WHITE, line=BORDER)
    rect(s, Inches(0.7), y, Pt(4), h, ORANGE if i < 2 else BLUE)
    text(s, Inches(0.9), y + Inches(0.08), Inches(3.1), Inches(0.42),
         [(k, 11.5, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.05), y + Inches(0.04), Inches(4.7), Inches(0.5),
         [(v, 11.5, GREY, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y += h + Inches(0.08)

picture(s, "digital_twin_box.png", Inches(9.25), Inches(2.30), Inches(3.35), Inches(2.2),
        "Digital Twin concept (illustrative)")
rect(s, Inches(9.25), Inches(4.85), Inches(3.35), Inches(1.6), LTBLUE, line=BORDER)
text(s, Inches(9.5), Inches(4.99), Inches(2.9), Inches(0.28),
     [("ORGANISATION", 10, BLUE, True)])
text(s, Inches(9.5), Inches(5.23), Inches(2.95), Inches(0.5),
     [("Ministry of Earth Sciences (MoES)", 12, NAVY, True)], line_spacing=1.12)
text(s, Inches(9.5), Inches(5.79), Inches(2.9), Inches(0.28),
     [("DEPARTMENT", 10, BLUE, True)])
text(s, Inches(9.5), Inches(6.03), Inches(2.95), Inches(0.4),
     [("India Meteorological Department", 12, NAVY, True)], line_spacing=1.12)

rect(s, Inches(0.7), Inches(6.37), Inches(8.2), Inches(0.52), LTORANGE, radius=True)
text(s, Inches(0.9), Inches(6.50), Inches(7.9), Inches(0.32),
     [("One portal for training, assessment and competency — that keeps working "
       "when the network does not.", 11.5, ORANGE, True)])

''' + s[b:]

# ---------------- 6. save: enforce the 6-slide limit ----------------
s = s.replace('prs.save(OUT)\nprint("saved", OUT)', '''# SIH allows a maximum of 6 slides - drop the template's instruction slide(s).
_ids = prs.slides._sldIdLst
_removed = 0
for _sld in list(_ids)[6:]:
    prs.part.drop_rel(_sld.rId)
    _ids.remove(_sld)
    _removed += 1

prs.save(OUT)
print("Filled your template ->", OUT)
if _removed:
    print("Removed %d extra slide(s) beyond the 6-slide limit." % _removed)
print("Slide size: %.2f x %.2f in (scale %.3f)"
      % (prs.slide_width / 914400, prs.slide_height / 914400, SCALE))
print("\\nNext: set Team Name / Team ID, check the layout, then export to PDF.")''')

open(DST, "w").write(s)
print("wrote", DST)
